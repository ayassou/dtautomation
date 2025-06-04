from crewai.tools import BaseTool
import PyPDF2
import docx
import pandas as pd
from pptx import Presentation  # Correction de l'importation
import os

class FileProcessingTool(BaseTool):
    name: str = "file_processing_tool"
    description: str = "Outil pour convertir des fichiers variés (PDF, Word, Excel, PPT) en texte brut et les découper en morceaux."

    def _run(self, file_paths: list) -> list:
        """
        Convertit des fichiers en texte brut et les découpe en morceaux.

        Args:
            file_paths (list): Liste des chemins vers les fichiers à traiter.

        Returns:
            list: Liste de dictionnaires contenant les morceaux de texte avec métadonnées.
        """
        processed_chunks = []

        for file_path in file_paths:
            if not os.path.exists(file_path):
                processed_chunks.append({"error": f"Fichier non trouvé : {file_path}"})
                continue

            file_name = os.path.basename(file_path)
            text = self._extract_text(file_path, file_name)

            if not text:
                processed_chunks.append({"error": f"Échec de l'extraction pour : {file_name}. Vérifiez le format ou l'intégrité du fichier."})
                continue

            chunks = self._chunk_text(text, file_name)
            processed_chunks.extend(chunks)

        return processed_chunks

    def _extract_text(self, file_path: str, file_name: str) -> str:
        """
        Extrait le texte brut d'un fichier selon son type.
        """
        if file_path.endswith(('.pdf', '.PDF')):
            try:
                with open(file_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    text = "".join(page.extract_text() for page in reader.pages if page.extract_text())
            except Exception as e:
                print(f"Erreur lors de l'extraction du PDF {file_name} : {str(e)}")
                text = ""
        elif file_path.endswith(('.docx', '.DOCX')):
            try:
                doc = docx.Document(file_path)
                text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
            except Exception as e:
                print(f"Erreur lors de l'extraction du DOCX {file_name} : {str(e)}")
                text = ""
        elif file_path.endswith(('.xlsx', '.xls', '.csv')):
            try:
                df = pd.read_excel(file_path) if file_path.endswith(('.xlsx', '.xls')) else pd.read_csv(file_path)
                text = "\n".join(" ".join(map(str, row)) for row in df.values)
            except Exception as e:
                print(f"Erreur lors de l'extraction du fichier Excel/CSV {file_name} : {str(e)}")
                text = ""
        elif file_path.endswith(('.pptx', '.PPTX')):
            try:
                ppt = Presentation(file_path)
                text = ""
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame") and shape.text_frame:
                            text += shape.text_frame.text + "\n"
            except Exception as e:
                print(f"Erreur lors de l'extraction du PPTX {file_name} : {str(e)}")
                text = ""
        else:
            text = ""
        return text.strip()

    def _chunk_text(self, text: str, file_name: str) -> list:
        """
        Découpe le texte en morceaux selon les règles : < 700 mots → tout prendre ; sinon 700 mots max avec chevauchement de 100 mots.
        Priorité 'Basse' pour les morceaux < 150 mots.
        """
        chunks = []
        words = text.split()
        total_words = len(words)

        if total_words < 700:
            chunk_text = " ".join(words)
            chunk_length = total_words
            priority = "Basse" if chunk_length < 150 else None
            chunks.append({
                "text": chunk_text,
                "part_id": 1,
                "priority": priority,
                "source": file_name
            })
        else:
            for i in range(0, total_words, 600):
                start_idx = max(0, i - 100)
                end_idx = min(i + 700, total_words)
                chunk_words = words[start_idx:end_idx]
                chunk_text = " ".join(chunk_words)
                chunk_length = len(chunk_words)

                priority = "Basse" if chunk_length < 150 else None

                chunks.append({
                    "text": chunk_text,
                    "part_id": len(chunks) + 1,
                    "priority": priority,
                    "source": file_name
                })

        return chunks

# Exemple d'utilisation (corrigé)
# tool = FileProcessingTool()
# # file_paths = ["Notice_Projet.pdf", "MGDIS - Aiden DATA.pptx"]
# file_paths = ["jira-issues-details.pdf"]
# result = tool._run(file_paths)
# for chunk in result:
#     if "error" in chunk:
#         print(f"Erreur : {chunk['error']}")
#     else:
#         print(f"Source: {chunk['source']}, Partie {chunk['part_id']}: {chunk['text'][:50]}... (Priorité: {chunk['priority']})")
